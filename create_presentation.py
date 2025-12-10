"""
Create PowerPoint presentation using python-pptx
Requires installation: pip install python-pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    """Create project presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    BLUE = RGBColor(31, 119, 180)
    DARK_BLUE = RGBColor(20, 60, 100)
    WHITE = RGBColor(255, 255, 255)
    
    # Slide 1: Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "MCP-Based RAG Multi-Agent\nResume Screening System"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DARK_BLUE
    
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.5))
    subtitle.text_frame.text = "Intelligent Resume Screening System Based on Retrieval-Augmented Generation"
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)
    subtitle.text_frame.paragraphs[0].font.color.rgb = BLUE
    
    # Slide 2: Project Background
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
    title = slide.shapes.title
    title.text = "Project Background and Motivation"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Problem Statement"
    p = tf.add_paragraph()
    p.text = "• HR departments receive hundreds of resumes daily"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Manual screening is time-consuming and may miss excellent candidates"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Lack of standardized evaluation criteria"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Solution"
    p = tf.add_paragraph()
    p.text = "• Automate resume screening process"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• AI-based intelligent matching analysis"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "• Structured evaluation reports"
    p.level = 1
    
    # Slide 3: Project Objectives
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Project Objectives"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. Efficiency Improvement"
    p = tf.add_paragraph()
    p.text = "   Automate batch processing of hundreds of resumes, reducing 80% of initial screening time"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "2. Accuracy Improvement"
    p = tf.add_paragraph()
    p.text = "   Semantic understanding-based intelligent matching, structured evaluation criteria"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "3. User Experience"
    p = tf.add_paragraph()
    p.text = "   Intuitive Web interface, clear visualization reports"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "4. Scalability"
    p = tf.add_paragraph()
    p.text = "   Support multiple LLM providers, local database easy to deploy"
    p.level = 1
    
    # Slide 4: System Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "System Architecture"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "User Layer: Streamlit UI"
    p = tf.add_paragraph()
    p.text = "Coordination Layer: MCP Client"
    p = tf.add_paragraph()
    p.text = "Tool Layer: MCP Server (PDF parsing/search)"
    p = tf.add_paragraph()
    p.text = "Data Layer: ChromaDB Vector Database"
    p = tf.add_paragraph()
    p.text = "AI Layer: Google Gemini LLM"
    
    # Save file
    prs.save('Resume_Screening_System_Presentation.pptx')
    print("Presentation created: Resume_Screening_System_Presentation.pptx")
    print("Note: This is a basic version, please refine content according to PRESENTATION_OUTLINE.md")

if __name__ == "__main__":
    try:
        create_presentation()
    except ImportError:
        print("Please install python-pptx first: pip install python-pptx")
    except Exception as e:
        print(f"Error: {e}")

